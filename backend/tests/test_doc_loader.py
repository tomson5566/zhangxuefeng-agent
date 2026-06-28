"""测试 doc_loader 5 格式解析。"""
import pytest
from backend.modules.doc_loader import save_upload, load_file, ALLOWED_EXT


def test_allowed_ext_set():
    assert ALLOWED_EXT == {'.txt', '.md', '.docx', '.xlsx', '.pdf', '.pptx'}


def test_load_txt(tmp_path):
    p = tmp_path / "test.txt"
    p.write_text("hello\nworld\n中文\n", encoding='utf-8')
    text = load_file(p)
    assert "hello" in text and "world" in text and "中文" in text


def test_load_docx(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("docx paragraph")
    doc.save(tmp_path / "test.docx")
    text = load_file(tmp_path / "test.docx")
    assert "docx paragraph" in text


def test_load_xlsx(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'name'
    ws['B1'] = 'score'
    ws['A2'] = 'alice'
    ws['B2'] = 95
    wb.save(tmp_path / "test.xlsx")
    text = load_file(tmp_path / "test.xlsx")
    assert 'alice' in text and '95' in text


def test_load_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX"
    prs.save(tmp_path / "test.pptx")
    text = load_file(tmp_path / "test.pptx")
    assert "Hello PPTX" in text


@pytest.mark.asyncio
async def test_save_upload_writes_files(tmp_path, monkeypatch):
    """save_upload 写 3 个文件:主 + .meta.json + .content.txt"""
    import backend.modules.doc_loader.loader as loader_mod
    monkeypatch.setattr(loader_mod, 'UPLOAD_ROOT', tmp_path)

    meta = await save_upload(b"hello content", "test.txt", "test-sid")
    assert meta['filename'] == 'test.txt'
    assert meta['size'] == 13
    assert meta['ext'] == '.txt'

    sid_dir = tmp_path / "test-sid"
    assert sid_dir.exists()
    files = list(sid_dir.iterdir())
    assert len(files) == 3
