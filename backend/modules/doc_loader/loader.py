"""文档上传 + 解析器。

支持格式:
- .txt / .md          → chardet 探测编码,直接读
- .docx               → python-docx 抽 paragraph + table 文本
- .xlsx               → openpyxl 抽所有 sheet 的 row 文本
- .pdf                → pypdf 抽每页 text
- .pptx               → python-pptx 抽每张 slide 的 shape 文本

公开 API:
- ALLOWED_EXT = {'.txt', '.md', '.docx', '.xlsx', '.pdf', '.pptx'}
- load_file(path: str | Path) -> str
- async save_upload(file_bytes, filename, session_id) -> dict
- list_uploaded(session_id) -> list[str]
"""
from __future__ import annotations
import json
import logging
import time
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)

ALLOWED_EXT: set[str] = {".txt", ".md", ".docx", ".xlsx", ".pdf", ".pptx"}
MAX_FILE_BYTES = 50 * 1024 * 1024  # 50MB
_PROJECT_ROOT = Path(__file__).resolve().parents[3]
UPLOAD_ROOT = _PROJECT_ROOT / "uploads"
UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)


def _detect_text(path: Path) -> str:
    import chardet
    raw = path.read_bytes()
    enc = chardet.detect(raw[:8192])["encoding"] or "utf-8"
    return raw.decode(enc, errors="replace")


def _load_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def _load_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _load_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"# Page {i + 1}\n{text}")
    return "\n".join(parts)


def _load_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        slide_text.append(p.text)
        if slide_text:
            parts.append(f"# Slide {i + 1}\n" + "\n".join(slide_text))
    return "\n".join(parts)


_LOADERS = {
    ".txt": _detect_text,
    ".md": _detect_text,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".pdf": _load_pdf,
    ".pptx": _load_pptx,
}


def load_file(path: str | Path) -> str:
    """按扩展名分发解析。返回纯文本。"""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Not found: {p}")
    ext = p.suffix.lower()
    if ext not in _LOADERS:
        raise ValueError(f"Unsupported ext: {ext}. Allowed: {sorted(ALLOWED_EXT)}")
    return _LOADERS[ext](p)


async def save_upload(file_bytes: bytes, filename: str, session_id: str) -> dict[str, Any]:
    """保存上传字节到 uploads/<sid>/<ts>_<safe>,返回 metadata。"""
    import aiofiles
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXT:
        raise ValueError(f"Unsupported ext: {ext}. Allowed: {sorted(ALLOWED_EXT)}")
    if len(file_bytes) > MAX_FILE_BYTES:
        raise ValueError(f"File too large: {len(file_bytes)} > {MAX_FILE_BYTES}")

    sid_dir = UPLOAD_ROOT / session_id
    sid_dir.mkdir(parents=True, exist_ok=True)

    safe = "".join(c if c.isalnum() or c in "._-" else "_" for c in filename)
    ts = int(time.time() * 1000)
    saved_path = sid_dir / f"{ts}_{safe}"
    async with aiofiles.open(saved_path, "wb") as f:
        await f.write(file_bytes)

    try:
        text = load_file(saved_path)
        preview = text[:500]
    except Exception as e:
        log.warning(f"Parse failed for {saved_path}: {e}")
        text = ""
        preview = f"[解析失败: {e}]"

    meta = {
        "path": str(saved_path),
        "filename": filename,
        "size": len(file_bytes),
        "ext": ext,
        "saved_at": ts,
        "session_id": session_id,
        "content_preview": preview,
        "content_full_len": len(text),
    }
    meta_path = saved_path.with_suffix(saved_path.suffix + ".meta.json")
    async with aiofiles.open(meta_path, "w", encoding="utf-8") as f:
        await f.write(json.dumps(meta, ensure_ascii=False, indent=2))

    content_path = saved_path.with_suffix(saved_path.suffix + ".content.txt")
    async with aiofiles.open(content_path, "w", encoding="utf-8") as f:
        await f.write(text)

    return meta


def list_uploaded(session_id: str) -> list[str]:
    """列已上传文件路径"""
    sid_dir = UPLOAD_ROOT / session_id
    if not sid_dir.exists():
        return []
    return sorted([str(p) for p in sid_dir.iterdir() if p.is_file()])


# 兼容旧 API(阶段 4 上一个版本的函数名)
def list_supported_extensions() -> list[str]:
    return sorted(ALLOWED_EXT)


# 别名 — 老代码用 SUPPORTED_EXTENSIONS,新代码用 ALLOWED_EXT
SUPPORTED_EXTENSIONS = ALLOWED_EXT
